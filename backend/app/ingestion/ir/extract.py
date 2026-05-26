"""Text extraction dispatcher — PDF / HTML / PPTX → plain text."""

import io
import logging
from typing import Literal

logger = logging.getLogger(__name__)

ExtractedKind = Literal["ir_pdf", "ir_html", "ir_pptx"]


def _extract_pdf(content: bytes) -> str:
    import pdfplumber

    text_parts: list[str] = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_html(content: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        text_parts.append(line)
    return "\n".join(text_parts)


def extract_text(content: bytes, content_type: str, url: str) -> tuple[str, ExtractedKind]:
    """Dispatch text extraction by content-type or URL extension.

    Returns (text, source_kind). The source_kind is one of the IR_* enum
    values stored in earnings_transcripts.source.
    """
    ctype = (content_type or "").lower()
    url_lower = url.lower()

    if "pdf" in ctype or url_lower.endswith(".pdf"):
        return _extract_pdf(content), "ir_pdf"
    if any(ext in ctype for ext in ("pptx", "powerpoint", "presentation")) or url_lower.endswith(".pptx"):
        return _extract_pptx(content), "ir_pptx"
    return _extract_html(content), "ir_html"
